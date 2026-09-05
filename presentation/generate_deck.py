#!/usr/bin/env python3
"""Generate the Hop OffGrid hackathon deck (16:9, seminar-bar house style)."""

from __future__ import annotations

import argparse
import sys
import uuid
import zipfile
from dataclasses import dataclass, field
from io import BytesIO
from pathlib import Path
from typing import Literal, Sequence
from xml.etree import ElementTree as ET

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
FONTS = ROOT / "fonts"
SHOTS = ROOT / "screenshots"
OUT_PPTX = ROOT / "Hop-OffGrid.pptx"
OUT_PDF = ROOT / "Hop-OffGrid.pdf"
REQUIRED_SHOTS = ("launch", "floor", "sheet", "blackout", "nearby")

# Exact 16:9 widescreen (13.333... in × 7.5 in)
SLIDE_W_EMU = 12_192_000
SLIDE_H_EMU = 6_858_000
SLIDE_W = SLIDE_W_EMU / 914_400
SLIDE_H = SLIDE_H_EMU / 914_400

# ~0.6 in margins
M = 0.60

# Locked type sizes (do not shrink per slide)
TITLE_PT = 32
BODY_PT = 16
META_PT = 11

SERIF = "Libre Baskerville"
SANS = "Inter"

# Surfaces
PAPER = (0xF5, 0xF3, 0xEF)
PANEL = (0xEE, 0xEC, 0xE8)
DARK = (0x12, 0x12, 0x12)
INK = (0x12, 0x12, 0x12)
MUTED = (0x5F, 0x5C, 0x56)
CREAM = (0xF5, 0xF3, 0xEF)
MUTED_DARK = (0xB0, 0xAD, 0xA6)

Face = Literal["serif", "serif-italic", "sans"]
Role = Literal["title", "body", "meta"]
Align = Literal["left", "right", "center"]


@dataclass
class Rect:
    x: float
    y: float
    w: float
    h: float
    fill: tuple[int, int, int]
    radius: float = 0.05


@dataclass
class TextBox:
    x: float
    y: float
    w: float
    h: float
    text: str
    face: Face
    role: Role
    color: tuple[int, int, int]
    align: Align = "left"
    valign: Literal["top", "middle"] = "top"
    tracking: int | None = None  # hundredths of a point


@dataclass
class ImageBox:
    x: float
    y: float
    w: float
    h: float
    path: Path


@dataclass
class SlideSpec:
    name: str
    bg: tuple[int, int, int]
    rects: list[Rect] = field(default_factory=list)
    texts: list[TextBox] = field(default_factory=list)
    images: list[ImageBox] = field(default_factory=list)


def rgb(c: tuple[int, int, int]) -> RGBColor:
    return RGBColor(*c)


def pt_for(role: Role) -> int:
    return {"title": TITLE_PT, "body": BODY_PT, "meta": META_PT}[role]


def font_file(face: Face) -> Path:
    return {
        "serif": FONTS / "LibreBaskerville-Regular.ttf",
        "serif-italic": FONTS / "LibreBaskerville-Italic.ttf",
        "sans": FONTS / "Inter-Regular.ttf",
    }[face]


def font_name(face: Face) -> str:
    return SERIF if face.startswith("serif") else SANS


# ---------------------------------------------------------------------------
# Layout helpers
# ---------------------------------------------------------------------------

def kicker_footer(n: int, total: int, kicker: str, dark: bool) -> list[TextBox]:
    mute = MUTED_DARK if dark else MUTED
    page = f"{n:02d} / {total:02d}"
    return [
        TextBox(M, M, 10.0, 0.26, kicker, "sans", "meta", mute, tracking=160),
        TextBox(M, 6.92, 7.4, 0.26, "AKSHAI KRISHNA S", "sans", "meta", mute, tracking=120),
        TextBox(5.93, 6.92, 6.80, 0.26, page, "sans", "meta", mute, align="right", tracking=120),
    ]


def paper_title(title: str) -> TextBox:
    return TextBox(M, 0.98, 12.13, 0.58, title, "serif", "title", INK)


def stack_panels(
    items: Sequence[tuple[str, str]],
    top: float = 1.74,
    height: float = 4.92,
    gap: float = 0.14,
) -> tuple[list[Rect], list[TextBox]]:
    n = len(items)
    ph = (height - gap * (n - 1)) / n
    width = SLIDE_W - 2 * M
    rects: list[Rect] = []
    texts: list[TextBox] = []
    pad = 0.22
    for i, (label, body) in enumerate(items):
        y = top + i * (ph + gap)
        rects.append(Rect(M, y, width, ph, PANEL))
        texts.append(TextBox(M + pad, y + 0.18, width - 2 * pad, 0.24, label, "serif-italic", "meta", MUTED))
        texts.append(TextBox(M + pad, y + 0.46, width - 2 * pad, ph - 0.64, body, "sans", "body", INK))
    return rects, texts


def grid_panels(
    items: Sequence[tuple[str, str]],
    cols: int,
    rows: int,
    top: float = 1.74,
    height: float = 4.92,
    gap: float = 0.16,
) -> tuple[list[Rect], list[TextBox]]:
    assert len(items) == cols * rows
    width = SLIDE_W - 2 * M
    cw = (width - gap * (cols - 1)) / cols
    ch = (height - gap * (rows - 1)) / rows
    pad = 0.22
    rects: list[Rect] = []
    texts: list[TextBox] = []
    for i, (label, body) in enumerate(items):
        c, r = i % cols, i // cols
        x = M + c * (cw + gap)
        y = top + r * (ch + gap)
        rects.append(Rect(x, y, cw, ch, PANEL))
        if label:
            texts.append(TextBox(x + pad, y + 0.20, cw - 2 * pad, 0.26, label, "serif-italic", "meta", MUTED))
            texts.append(TextBox(x + pad, y + 0.52, cw - 2 * pad, ch - 0.74, body, "sans", "body", INK))
        else:
            texts.append(
                TextBox(x + pad, y + 0.28, cw - 2 * pad, ch - 0.56, body, "sans", "body", INK, valign="middle")
            )
    return rects, texts


def shot(name: str) -> Path:
    path = SHOTS / f"{name}.png"
    if not path.exists():
        raise FileNotFoundError(f"Missing screenshot: {path}")
    return path


def phone_aspect(path: Path) -> float:
    from PIL import Image

    with Image.open(path) as im:
        w, h = im.size
    return w / h


def phone_row(
    items: Sequence[tuple[str, str]],
    top: float = 1.74,
    height: float = 4.96,
    gap: float = 0.36,
) -> tuple[list[ImageBox], list[TextBox]]:
    """Centered row of real device shots with italic labels. No frames."""
    label_h = 0.30
    phone_h = height - label_h
    images: list[ImageBox] = []
    texts: list[TextBox] = []
    sized: list[tuple[Path, str, float, float]] = []
    for name, label in items:
        path = shot(name)
        pw = phone_h * phone_aspect(path)
        sized.append((path, label, pw, phone_h))
    total_w = sum(pw for _, _, pw, _ in sized) + gap * (len(sized) - 1)
    x = (SLIDE_W - total_w) / 2
    for path, label, pw, ph in sized:
        images.append(ImageBox(x, top, pw, ph, path))
        texts.append(
            TextBox(x - 0.08, top + ph + 0.04, pw + 0.16, 0.26, label, "serif-italic", "meta", MUTED, align="center")
        )
        x += pw + gap
    return images, texts


def two_col_lists(
    left: Sequence[tuple[str, str]],
    right: Sequence[tuple[str, str]],
    top: float = 1.74,
    height: float = 4.92,
    gap: float = 0.16,
) -> tuple[list[Rect], list[TextBox]]:
    width = SLIDE_W - 2 * M
    cw = (width - gap) / 2
    pad_x = 0.26
    pad_y = 0.24
    rects = [
        Rect(M, top, cw, height, PANEL),
        Rect(M + cw + gap, top, cw, height, PANEL),
    ]
    texts: list[TextBox] = []

    def fill(x: float, items: Sequence[tuple[str, str]]) -> None:
        inner_h = height - 2 * pad_y
        row_h = inner_h / len(items)
        for i, (label, body) in enumerate(items):
            y = top + pad_y + i * row_h
            texts.append(TextBox(x + pad_x, y + 0.04, cw - 2 * pad_x, 0.22, label, "serif-italic", "meta", MUTED))
            texts.append(TextBox(x + pad_x, y + 0.28, cw - 2 * pad_x, row_h - 0.36, body, "sans", "body", INK))

    fill(M, left)
    fill(M + cw + gap, right)
    return rects, texts


def build_slides() -> list[SlideSpec]:
    total = 9
    slides: list[SlideSpec] = []

    # 1. Title
    t1 = SlideSpec("title", DARK, texts=kicker_footer(1, total, "01  ·  OFFGRID", dark=True))
    t1.texts += [
        TextBox(M, 2.28, 12.13, 0.58, "Hop", "serif", "title", CREAM),
        TextBox(M, 2.96, 12.13, 0.40, "Your floor, offline.", "serif-italic", "body", CREAM),
        TextBox(M, 3.92, 12.13, 0.30, "OffGrid  ·  solo  ·  open source", "sans", "body", MUTED_DARK),
        TextBox(M, 4.30, 12.13, 0.30, "github.com/flyme2mars/hop", "sans", "body", MUTED_DARK),
    ]
    slides.append(t1)

    # 2. Problem
    t2 = SlideSpec("problem", PAPER, texts=kicker_footer(2, total, "02  ·  PROBLEM", dark=False))
    t2.texts.append(paper_title("Hostel floors lose the network."))
    r, tx = stack_panels(
        [
            ("Dead zones", "Hostel floors often have dead zones, or no usable network."),
            ("Power cuts", "During a power cut, phones cannot help the floor coordinate."),
            ("The tower", "WhatsApp needs the tower. The floor does not always have it."),
        ]
    )
    t2.rects += r
    t2.texts += tx
    slides.append(t2)

    # 3. Solution
    t3 = SlideSpec("solution", PAPER, texts=kicker_footer(3, total, "03  ·  SOLUTION", dark=False))
    t3.texts.append(paper_title("A floor board with no internet."))
    r, tx = stack_panels(
        [
            (
                "Android app",
                "A local Offer / Ask / Note board on your floor. It works with no internet.",
            ),
            (
                "Nearby phones",
                "Phones on the same floor can share posts over Bluetooth when the app is open.",
            ),
            (
                "Blackout mode",
                "A full black screen with a timer, and I’m OK / Need help.",
            ),
        ]
    )
    t3.rects += r
    t3.texts += tx
    slides.append(t3)

    # 4. Features
    t4 = SlideSpec("features", PAPER, texts=kicker_footer(4, total, "04  ·  FEATURES", dark=False))
    t4.texts.append(paper_title("Setup, board, blackout, nearby sync."))
    r, tx = two_col_lists(
        [
            ("Setup", "Name, room, floor."),
            ("Floor board", "Board with filters."),
            ("Claim", "Claim posts."),
            ("History", "History of posts."),
        ],
        [
            ("Settings", "Settings."),
            ("Blackout mode", "Full black screen with a timer and I’m OK / Need help."),
            ("Bluetooth nearby sync", "Same floor, when the app is open."),
            ("Offline alone", "Fully usable offline alone."),
        ],
    )
    t4.rects += r
    t4.texts += tx
    slides.append(t4)

    # 5. Screens — first open, floor, new post
    t5 = SlideSpec("board", PAPER, texts=kicker_footer(5, total, "05  ·  SCREENS", dark=False))
    t5.texts.append(paper_title("The floor board."))
    imgs, labels = phone_row(
        [
            ("launch", "First open"),
            ("floor", "Floor"),
            ("sheet", "New post"),
        ],
        gap=0.34,
    )
    t5.images += imgs
    t5.texts += labels
    slides.append(t5)

    # 6. Screens — blackout and nearby
    t6 = SlideSpec("states", PAPER, texts=kicker_footer(6, total, "06  ·  SCREENS", dark=False))
    t6.texts.append(paper_title("Blackout, and nearby."))
    imgs, labels = phone_row(
        [
            ("blackout", "Blackout"),
            ("nearby", "Nearby"),
        ],
        gap=0.55,
    )
    t6.images += imgs
    t6.texts += labels
    slides.append(t6)

    # 7. Tech
    t7 = SlideSpec("stack", PAPER, texts=kicker_footer(7, total, "07  ·  STACK", dark=False))
    t7.texts.append(paper_title("What it is built with."))
    r, tx = grid_panels(
        [
            ("App", "Kotlin. Jetpack Compose. Material 3."),
            ("Local store", "Room. DataStore."),
            ("Nearby", "BLE advertise / scan + GATT sync."),
            ("Build", "GitHub Actions debug APK."),
        ],
        cols=2,
        rows=2,
    )
    t7.rects += r
    t7.texts += tx
    slides.append(t7)

    # 8. Future
    t8 = SlideSpec("future", PAPER, texts=kicker_footer(8, total, "08  ·  FUTURE", dark=False))
    t8.texts.append(paper_title("After the first floor."))
    r, tx = grid_panels(
        [
            ("Mesh", "Stronger multi-hop mesh."),
            ("Peers", "Better peer UX."),
            ("Background", "Background limits / battery."),
            ("iOS", "iOS if ever needed."),
        ],
        cols=2,
        rows=2,
    )
    t8.rects += r
    t8.texts += tx
    slides.append(t8)

    # 9. Close
    t9 = SlideSpec("close", DARK, texts=kicker_footer(9, total, "09  ·  OFFGRID", dark=True))
    t9.texts += [
        TextBox(M, 2.28, 12.13, 0.58, "Hop", "serif", "title", CREAM),
        TextBox(M, 2.96, 12.13, 0.40, "Your floor, offline.", "serif-italic", "body", CREAM),
        TextBox(M, 3.92, 12.13, 0.30, "github.com/flyme2mars/hop", "sans", "body", MUTED_DARK),
        TextBox(M, 4.30, 12.13, 0.30, "Solo OffGrid entry  ·  open source", "sans", "body", MUTED_DARK),
    ]
    slides.append(t9)

    return slides


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------

def _set_run(run, face: Face, role: Role, color: tuple[int, int, int], tracking: int | None) -> None:
    run.font.name = font_name(face)
    run.font.size = Pt(pt_for(role))
    run.font.bold = False
    run.font.italic = face == "serif-italic"
    run.font.color.rgb = rgb(color)
    rPr = run._r.get_or_add_rPr()
    if tracking is not None:
        rPr.set("spc", str(tracking))


def add_textbox(slide, box: TextBox) -> None:
    shape = slide.shapes.add_textbox(Inches(box.x), Inches(box.y), Inches(box.w), Inches(box.h))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.anchor = MSO_ANCHOR.MIDDLE if box.valign == "middle" else MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "right": PP_ALIGN.RIGHT, "center": PP_ALIGN.CENTER}[box.align]
    p.clear()
    run = p.add_run()
    run.text = box.text
    _set_run(run, box.face, box.role, box.color, box.tracking)

    pPr = p._p.get_or_add_pPr()
    ln = etree.SubElement(pPr, qn("a:lnSpc"))
    spc_pct = etree.SubElement(ln, qn("a:spcPct"))
    spc_pct.set("val", "120000" if box.role == "title" else "135000" if box.role == "body" else "120000")


def add_rect(slide, rect: Rect) -> None:
    adj = 0 if rect.radius <= 0 else min(rect.radius, min(rect.w, rect.h) / 2)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if adj else MSO_SHAPE.RECTANGLE,
        Inches(rect.x),
        Inches(rect.y),
        Inches(rect.w),
        Inches(rect.h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(rect.fill)
    shape.line.fill.background()
    # No shadow
    spPr = shape._element.spPr
    effect = spPr.find(qn("a:effectLst"))
    if effect is not None:
        spPr.remove(effect)
    if adj and shape.adjustments:
        # adjustment 0 is corner radius as fraction of half-min-side
        half = min(rect.w, rect.h) / 2
        shape.adjustments[0] = adj / half if half else 0


def set_slide_bg(slide, color: tuple[int, int, int]) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def write_pptx(slides: Sequence[SlideSpec], path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W_EMU)
    prs.slide_height = Emu(SLIDE_H_EMU)
    blank = prs.slide_layouts[6]
    for spec in slides:
        slide = prs.slides.add_slide(blank)
        set_slide_bg(slide, spec.bg)
        for rect in spec.rects:
            add_rect(slide, rect)
        for img in spec.images:
            slide.shapes.add_picture(str(img.path), Inches(img.x), Inches(img.y), Inches(img.w), Inches(img.h))
        for box in spec.texts:
            add_textbox(slide, box)
    prs.core_properties.title = "Hop · OffGrid"
    prs.core_properties.author = "Akshai Krishna S"
    prs.core_properties.subject = "OffGrid hackathon presentation"
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    _patch_theme_fonts(path)
    _embed_fonts(path)


def _patch_theme_fonts(path: Path) -> None:
    """Point theme major/minor fonts at the bundled families."""
    buf = BytesIO()
    with zipfile.ZipFile(path, "r") as zin, zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.startswith("ppt/theme/theme") and item.filename.endswith(".xml"):
                text = data.decode("utf-8")
                text = text.replace('typeface="Calibri Light"', f'typeface="{SERIF}"')
                text = text.replace('typeface="Calibri"', f'typeface="{SANS}"')
                text = text.replace('typeface="+mj-lt"', f'typeface="{SERIF}"')
                data = text.encode("utf-8")
            zout.writestr(item, data)
    path.write_bytes(buf.getvalue())


def _guid_bytes(u: uuid.UUID) -> bytes:
    """Microsoft GUID mixed-endian 16-byte layout."""
    return u.bytes_le


def _obfuscate_font(data: bytes, guid: uuid.UUID) -> bytes:
    key = _guid_bytes(guid) * 2
    out = bytearray(data)
    for i in range(min(32, len(out))):
        out[i] ^= key[i]
    return bytes(out)


def _embed_fonts(path: Path) -> None:
    """Embed Inter + Libre Baskerville as OOXML font parts (obfuscated TTF)."""
    fonts = [
        ("Inter", "regular", FONTS / "Inter-Regular.ttf"),
        ("Inter", "italic", FONTS / "Inter-Italic.ttf"),
        ("Libre Baskerville", "regular", FONTS / "LibreBaskerville-Regular.ttf"),
        ("Libre Baskerville", "italic", FONTS / "LibreBaskerville-Italic.ttf"),
        ("Libre Baskerville", "bold", FONTS / "LibreBaskerville-Bold.ttf"),
    ]
    P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    CT = "http://schemas.openxmlformats.org/package/2006/content-types"
    REL_FONT = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/font"

    ET.register_namespace("p", P)
    ET.register_namespace("r", R)
    ET.register_namespace("", CT)

    with zipfile.ZipFile(path, "r") as zin:
        names = set(zin.namelist())
        files = {n: zin.read(n) for n in names}

    pres_xml = ET.fromstring(files["ppt/presentation.xml"])
    rels_xml = ET.fromstring(files["ppt/_rels/presentation.xml.rels"])
    ct_xml = ET.fromstring(files["[Content_Types].xml"])

    # Existing rIds
    existing = []
    for rel in rels_xml:
        rid = rel.attrib.get("Id", "")
        if rid.startswith("rId"):
            try:
                existing.append(int(rid[3:]))
            except ValueError:
                pass
    next_id = max(existing, default=1) + 1

    # Content type
    if not any(el.attrib.get("Extension") == "fntdata" for el in ct_xml):
        ET.SubElement(
            ct_xml,
            f"{{{CT}}}Default",
            {"Extension": "fntdata", "ContentType": "application/x-fontdata"},
        )

    # Group by typeface
    by_face: dict[str, dict[str, tuple[int, bytes, uuid.UUID]]] = {}
    font_parts: dict[str, bytes] = {}
    for typeface, style, fpath in fonts:
        guid = uuid.uuid4()
        rid_n = next_id
        next_id += 1
        part_name = f"ppt/fonts/font{rid_n}.fntdata"
        font_parts[part_name] = _obfuscate_font(fpath.read_bytes(), guid)
        rel = ET.SubElement(
            rels_xml,
            f"{{{R}}}Relationship",
            {
                "Id": f"rId{rid_n}",
                "Type": REL_FONT,
                "Target": f"fonts/font{rid_n}.fntdata",
            },
        )
        by_face.setdefault(typeface, {})[style] = (rid_n, font_parts[part_name], guid)

    # Remove any existing embeddedFontLst
    for child in list(pres_xml):
        if child.tag == f"{{{P}}}embeddedFontLst":
            pres_xml.remove(child)

    lst = ET.Element(f"{{{P}}}embeddedFontLst")
    # Insert after p:sldIdLst if present, else append
    insert_at = 0
    for i, child in enumerate(list(pres_xml)):
        if child.tag.endswith("sldIdLst"):
            insert_at = i + 1
    pres_xml.insert(insert_at, lst)

    style_tag = {
        "regular": "regular",
        "italic": "italic",
        "bold": "bold",
        "boldItalic": "boldItalic",
    }
    for typeface, styles in by_face.items():
        node = ET.SubElement(lst, f"{{{P}}}embeddedFont")
        ET.SubElement(node, f"{{{P}}}font", {"typeface": typeface, "pitchFamily": "2", "charset": "0"})
        for style, (rid_n, _, _) in styles.items():
            ET.SubElement(node, f"{{{P}}}{style_tag[style]}", {f"{{{R}}}id": f"rId{rid_n}"})

    files["ppt/presentation.xml"] = ET.tostring(pres_xml, encoding="utf-8", xml_declaration=True)
    files["ppt/_rels/presentation.xml.rels"] = ET.tostring(rels_xml, encoding="utf-8", xml_declaration=True)
    files["[Content_Types].xml"] = ET.tostring(ct_xml, encoding="utf-8", xml_declaration=True)
    files.update(font_parts)

    buf = BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zout:
        for name, data in files.items():
            zout.writestr(name, data)
    path.write_bytes(buf.getvalue())


# ---------------------------------------------------------------------------
# Visual QA (PIL) + geometry
# ---------------------------------------------------------------------------

def _wrap_lines(text: str, font, max_width_px: float) -> list[str]:
    if not text:
        return [""]
    words = text.split()
    lines: list[str] = []
    cur: list[str] = []
    for word in words:
        trial = (" ".join(cur + [word])).strip()
        if font.getlength(trial) <= max_width_px or not cur:
            cur.append(word)
        else:
            lines.append(" ".join(cur))
            cur = [word]
    if cur:
        lines.append(" ".join(cur))
    return lines or [""]


def render_qa(slides: Sequence[SlideSpec], out_dir: Path, dpi: int = 144) -> list[Path]:
    from PIL import Image, ImageDraw, ImageFont

    out_dir.mkdir(parents=True, exist_ok=True)
    px_w, px_h = int(SLIDE_W * dpi), int(SLIDE_H * dpi)
    cache: dict[tuple[Face, Role], ImageFont.FreeTypeFont] = {}

    def font(face: Face, role: Role) -> ImageFont.FreeTypeFont:
        key = (face, role)
        if key not in cache:
            # pt → px at this dpi
            px = int(round(pt_for(role) * dpi / 72))
            cache[key] = ImageFont.truetype(str(font_file(face)), px)
        return cache[key]

    issues: list[str] = []
    paths: list[Path] = []

    def inch_box(b) -> tuple[float, float, float, float]:
        return b.x, b.y, b.x + b.w, b.y + b.h

    def overlap(a, b, eps: float = 0.02) -> bool:
        ax1, ay1, ax2, ay2 = a
        bx1, by1, bx2, by2 = b
        return ax1 < bx2 - eps and ax2 > bx1 + eps and ay1 < by2 - eps and ay2 > by1 + eps

    for i, spec in enumerate(slides, start=1):
        im = Image.new("RGB", (px_w, px_h), spec.bg)
        dr = ImageDraw.Draw(im)

        # Geometry: texts must stay inside the slide; panels must not overlap each other
        for a in spec.rects:
            if a.x < -0.01 or a.y < -0.01 or a.x + a.w > SLIDE_W + 0.01 or a.y + a.h > SLIDE_H + 0.01:
                issues.append(f"slide {i} panel outside slide: {a}")
        for a, b in (
            (spec.rects[j], spec.rects[k])
            for j in range(len(spec.rects))
            for k in range(j + 1, len(spec.rects))
        ):
            if overlap(inch_box(a), inch_box(b)):
                issues.append(f"slide {i} overlapping panels")
        for img in spec.images:
            if img.x < -0.01 or img.y < -0.01 or img.x + img.w > SLIDE_W + 0.01 or img.y + img.h > SLIDE_H + 0.01:
                issues.append(f"slide {i} image outside slide: {img.path.name}")
        for a, b in (
            (spec.images[j], spec.images[k])
            for j in range(len(spec.images))
            for k in range(j + 1, len(spec.images))
        ):
            if overlap(inch_box(a), inch_box(b)):
                issues.append(f"slide {i} overlapping shots")

        for rect in spec.rects:
            x1, y1 = rect.x * dpi, rect.y * dpi
            x2, y2 = (rect.x + rect.w) * dpi, (rect.y + rect.h) * dpi
            rad = rect.radius * dpi
            dr.rounded_rectangle((x1, y1, x2, y2), radius=rad, fill=rect.fill)

        for img in spec.images:
            src = Image.open(img.path).convert("RGBA")
            dest_size = (max(1, int(img.w * dpi)), max(1, int(img.h * dpi)))
            fitted = src.resize(dest_size, Image.Resampling.LANCZOS)
            paper = Image.new("RGBA", dest_size, spec.bg + (255,))
            paper.alpha_composite(fitted)
            im.paste(paper.convert("RGB"), (int(img.x * dpi), int(img.y * dpi)))

        for box in spec.texts:
            fnt = font(box.face, box.role)
            max_w = box.w * dpi
            lines = _wrap_lines(box.text, fnt, max_w)
            line_gap = pt_for(box.role) * dpi / 72 * (1.20 if box.role != "body" else 1.35)
            block_h = line_gap * len(lines)
            if box.valign == "middle":
                y0 = box.y * dpi + (box.h * dpi - block_h) / 2
            else:
                y0 = box.y * dpi
            if block_h > box.h * dpi + 2:
                issues.append(
                    f"slide {i} overflow: {box.text!r} needs {block_h/dpi:.2f}in in {box.h:.2f}in box"
                )
            if box.x < -0.01 or box.y < -0.01 or box.x + box.w > SLIDE_W + 0.01:
                issues.append(f"slide {i} text box outside slide: {box.text!r}")
            for li, line in enumerate(lines):
                lw = fnt.getlength(line)
                if box.align == "right":
                    x = (box.x + box.w) * dpi - lw
                elif box.align == "center":
                    x = box.x * dpi + (max_w - lw) / 2
                else:
                    x = box.x * dpi
                dr.text((x, y0 + li * line_gap), line, font=fnt, fill=box.color)

        dest = out_dir / f"slide-{i:02d}-{spec.name}.png"
        im.save(dest, "PNG")
        paths.append(dest)

    if issues:
        print("VISUAL QA FAILED:", file=sys.stderr)
        for msg in issues:
            print(" -", msg, file=sys.stderr)
        raise SystemExit(2)
    print(f"Visual QA clean: {len(slides)} slides, no overflow, no panel overlap.")
    return paths


# ---------------------------------------------------------------------------
# PDF (reportlab) — optional twin with the same boxes
# ---------------------------------------------------------------------------

def write_pdf(slides: Sequence[SlideSpec], path: Path) -> None:
    from reportlab.lib.colors import Color
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.pdfgen import canvas

    pdfmetrics.registerFont(TTFont("LibreBaskerville", str(FONTS / "LibreBaskerville-Regular.ttf")))
    pdfmetrics.registerFont(TTFont("LibreBaskerville-Italic", str(FONTS / "LibreBaskerville-Italic.ttf")))
    pdfmetrics.registerFont(TTFont("Inter", str(FONTS / "Inter-Regular.ttf")))

    def rl_font(face: Face) -> str:
        return {
            "serif": "LibreBaskerville",
            "serif-italic": "LibreBaskerville-Italic",
            "sans": "Inter",
        }[face]

    def rl_color(c: tuple[int, int, int]) -> Color:
        return Color(c[0] / 255, c[1] / 255, c[2] / 255)

    W, H = SLIDE_W * 72, SLIDE_H * 72
    c = canvas.Canvas(str(path), pagesize=(W, H))
    c.setTitle("Hop · OffGrid")
    c.setAuthor("Akshai Krishna S")

    def y_top(inch_y: float, inch_h: float) -> float:
        return H - (inch_y + inch_h) * 72

    for spec in slides:
        c.setFillColor(rl_color(spec.bg))
        c.rect(0, 0, W, H, fill=1, stroke=0)
        for rect in spec.rects:
            c.setFillColor(rl_color(rect.fill))
            x, y = rect.x * 72, y_top(rect.y, rect.h)
            rad = rect.radius * 72
            c.roundRect(x, y, rect.w * 72, rect.h * 72, rad, fill=1, stroke=0)
        for img in spec.images:
            c.drawImage(
                str(img.path),
                img.x * 72,
                y_top(img.y, img.h),
                width=img.w * 72,
                height=img.h * 72,
                preserveAspectRatio=True,
                mask="auto",
            )
        for box in spec.texts:
            c.setFillColor(rl_color(box.color))
            c.setFont(rl_font(box.face), pt_for(box.role))
            # wrap
            from reportlab.pdfbase.pdfmetrics import stringWidth

            max_w = box.w * 72
            words = box.text.split()
            lines: list[str] = []
            cur: list[str] = []
            for word in words:
                trial = " ".join(cur + [word])
                if stringWidth(trial, rl_font(box.face), pt_for(box.role)) <= max_w or not cur:
                    cur.append(word)
                else:
                    lines.append(" ".join(cur))
                    cur = [word]
            if cur:
                lines.append(" ".join(cur))
            leading = pt_for(box.role) * (1.20 if box.role != "body" else 1.35)
            block_h = leading * len(lines)
            if box.valign == "middle":
                top = H - box.y * 72 - (box.h * 72 - block_h) / 2
            else:
                top = H - box.y * 72
            for li, line in enumerate(lines):
                lw = stringWidth(line, rl_font(box.face), pt_for(box.role))
                if box.align == "right":
                    x = (box.x + box.w) * 72 - lw
                elif box.align == "center":
                    x = box.x * 72 + (max_w - lw) / 2
                else:
                    x = box.x * 72
                # reportlab drawString baseline ≈ 0.8em below top of line box
                baseline = top - (li * leading) - pt_for(box.role) * 0.85
                extra = {}
                if box.tracking:
                    extra["charSpace"] = box.tracking / 100.0
                c.drawString(x, baseline, line, **extra)
        c.showPage()
    c.save()


def inspect_pptx(path: Path) -> int:
    prs = Presentation(str(path))
    w_in = prs.slide_width / 914_400
    h_in = prs.slide_height / 914_400
    n = len(prs.slides)
    print(f"PPTX {path.name}: {n} slides, {w_in:.3f} in × {h_in:.3f} in")
    if n > 10:
        raise SystemExit(f"Too many slides: {n}")
    if abs(w_in - 13.333) > 0.02 or abs(h_in - 7.5) > 0.02:
        raise SystemExit(f"Not 16:9 13.333×7.5: {w_in}×{h_in}")
    return n


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--qa-dir", default=str(ROOT / "qa"), help="PNG preview directory")
    parser.add_argument("--no-pdf", action="store_true")
    parser.add_argument("--no-qa", action="store_true")
    args = parser.parse_args()

    for needed in (
        FONTS / "Inter-Regular.ttf",
        FONTS / "LibreBaskerville-Regular.ttf",
        FONTS / "LibreBaskerville-Italic.ttf",
    ):
        if not needed.exists():
            raise SystemExit(f"Missing font: {needed}")
    for name in REQUIRED_SHOTS:
        if not shot(name).exists():
            raise SystemExit(f"Missing screenshot: {shot(name)}")

    slides = build_slides()
    if len(slides) > 10:
        raise SystemExit("Deck exceeds 10 slides")

    write_pptx(slides, OUT_PPTX)
    n = inspect_pptx(OUT_PPTX)
    print(f"Wrote {OUT_PPTX} ({n} slides)")

    if not args.no_qa:
        render_qa(slides, Path(args.qa_dir))

    if not args.no_pdf:
        try:
            write_pdf(slides, OUT_PDF)
            print(f"Wrote {OUT_PDF}")
        except ImportError:
            print("reportlab not installed; skip PDF", file=sys.stderr)


if __name__ == "__main__":
    main()
